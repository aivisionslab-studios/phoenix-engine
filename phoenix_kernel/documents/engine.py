"""
phoenix_kernel/documents/engine.py

Extração e reconstrução de documentos (PDF, DOCX, XLSX, PPTX, TXT/MD) para
a Document Engine da Phoenix.

Mantém o mesmo princípio dos outros drivers (vision, piper, sd_cpp): o
ResidentManager não sabe NADA sobre bibliotecas de documento - só chama
extract_text()/rebuild_document() e manda o texto pro LLM via
runtime.execute(), exatamente como já faz com imagem (sd_cpp.py) e voz
(piper.py). Esse módulo é o único lugar que conhece pymupdf4llm/
python-docx/openpyxl/python-pptx.

Todas as funções aqui são SÍNCRONAS de propósito (I/O + CPU puro, sem
GPU) - quem chamar deve rodar via loop.run_in_executor() para não
bloquear o event loop asyncio, igual reasoning_engine.py já faz com
self.knowledge.build_context.
"""

import logging
from pathlib import Path

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx", ".txt", ".md"}

# PDF nunca é usado como formato de SAÍDA (a extração de texto perde toda
# a formatação visual, então "reconstruir um PDF" a partir de texto puro
# seria uma mentira de qualidade - por isso REBUILD_EXTENSIONS não inclui
# .pdf; edições de PDF sempre saem como .docx, que é editável de verdade).
REBUILD_EXTENSIONS = {".docx", ".xlsx", ".pptx", ".txt", ".md"}


class DocumentEngineError(Exception):
    """Erro ao extrair ou reconstruir um documento. Sempre carrega uma
    mensagem clara o suficiente para virar a resposta HTTP de erro."""


# ============================================================
# EXTRAÇÃO
# ============================================================

def extract_text(path: Path) -> str:
    """Extrai o conteúdo textual/estruturado de um documento como Markdown
    (PDF/DOCX/PPTX) ou texto tabular (XLSX) ou texto puro (TXT/MD).

    Lança DocumentEngineError com uma mensagem clara se a extensão não for
    suportada, a dependência não estiver instalada, ou a extração falhar -
    nunca devolve string vazia silenciosamente (isso já causou bug real
    no endpoint antigo /api/documents/ingest, que devolvia
    {"error": "..."} só depois de já ter tentado mandar texto vazio pro
    LLM).
    """
    ext = path.suffix.lower()

    if ext == ".pdf":
        return _extract_pdf(path)
    if ext == ".docx":
        return _extract_docx(path)
    if ext == ".xlsx":
        return _extract_xlsx(path)
    if ext == ".pptx":
        return _extract_pptx(path)
    if ext in (".txt", ".md"):
        return _extract_plain_text(path)

    raise DocumentEngineError(
        f"Extensão '{ext}' não suportada. Formatos aceitos: "
        f"{', '.join(sorted(SUPPORTED_EXTENSIONS))}."
    )


def _extract_pdf(path: Path) -> str:
    text = ""
    try:
        import pymupdf4llm
        text = pymupdf4llm.to_markdown(str(path))
    except ImportError:
        try:
            import fitz  # PyMuPDF puro - fallback se pymupdf4llm não estiver instalado
        except ImportError as e:
            raise DocumentEngineError(
                "Nenhuma biblioteca de PDF instalada (pymupdf4llm/pymupdf). "
                "Rode: pip install pymupdf pymupdf4llm"
            ) from e
        doc = fitz.open(str(path))
        try:
            text = "\n\n".join(page.get_text() for page in doc)
        finally:
            doc.close()
    except Exception as e:
        raise DocumentEngineError(f"Falha ao extrair texto do PDF: {e}") from e

    if not text or not text.strip():
        raise DocumentEngineError(
            "PDF extraído está vazio (pode ser um PDF escaneado sem OCR - "
            "esta versão da Document Engine não faz OCR)."
        )
    return text


def _extract_docx(path: Path) -> str:
    try:
        import docx
    except ImportError as e:
        raise DocumentEngineError("python-docx não instalado. Rode: pip install python-docx") from e

    try:
        document = docx.Document(str(path))
    except Exception as e:
        raise DocumentEngineError(f"Falha ao abrir o DOCX: {e}") from e

    parts = []
    for para in document.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))

    text = "\n\n".join(parts)
    if not text.strip():
        raise DocumentEngineError("DOCX extraído está vazio.")
    return text


def _extract_xlsx(path: Path) -> str:
    try:
        import openpyxl
    except ImportError as e:
        raise DocumentEngineError("openpyxl não instalado. Rode: pip install openpyxl") from e

    try:
        wb = openpyxl.load_workbook(str(path), data_only=True)
    except Exception as e:
        raise DocumentEngineError(f"Falha ao abrir o XLSX: {e}") from e

    parts = []
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        parts.append(f"## Planilha: {sheet_name}")
        for row in sheet.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                parts.append(" | ".join("" if c is None else str(c) for c in row))

    text = "\n".join(parts)
    if not text.strip():
        raise DocumentEngineError("XLSX extraído está vazio.")
    return text


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as e:
        raise DocumentEngineError("python-pptx não instalado. Rode: pip install python-pptx") from e

    try:
        prs = Presentation(str(path))
    except Exception as e:
        raise DocumentEngineError(f"Falha ao abrir o PPTX: {e}") from e

    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_lines = [f"## Slide {i}"]
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                slide_lines.append(shape.text_frame.text)
        if len(slide_lines) > 1:
            parts.append("\n".join(slide_lines))

    text = "\n\n".join(parts)
    if not text.strip():
        raise DocumentEngineError("PPTX extraído está vazio.")
    return text


def _extract_plain_text(path: Path) -> str:
    try:
        text = path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        text = path.read_text(encoding="latin-1")
    except Exception as e:
        raise DocumentEngineError(f"Falha ao ler arquivo de texto: {e}") from e

    if not text.strip():
        raise DocumentEngineError("Arquivo de texto está vazio.")
    return text


# ============================================================
# RECONSTRUÇÃO — recebe o texto (já editado pelo LLM) e grava um novo
# arquivo. A extensão de `output_path` decide o formato de saída.
# ============================================================

def rebuild_document(new_text: str, output_path: Path) -> Path:
    """Grava `new_text` em `output_path`, escolhendo o formato pela
    extensão de `output_path`. Lança DocumentEngineError se a extensão
    de saída não for suportada (ex: .pdf - ver REBUILD_EXTENSIONS)."""
    ext = output_path.suffix.lower()
    if ext not in REBUILD_EXTENSIONS:
        raise DocumentEngineError(
            f"Não é possível gerar saída no formato '{ext}'. "
            f"Formatos de saída suportados: {', '.join(sorted(REBUILD_EXTENSIONS))}."
        )

    output_path.parent.mkdir(parents=True, exist_ok=True)

    if ext == ".docx":
        _rebuild_docx(new_text, output_path)
    elif ext == ".xlsx":
        _rebuild_xlsx(new_text, output_path)
    elif ext == ".pptx":
        _rebuild_pptx(new_text, output_path)
    else:  # .txt / .md
        output_path.write_text(new_text, encoding="utf-8")

    return output_path


def _rebuild_docx(text: str, output_path: Path) -> None:
    try:
        import docx
    except ImportError as e:
        raise DocumentEngineError("python-docx não instalado. Rode: pip install python-docx") from e

    document = docx.Document()
    for line in text.split("\n"):
        document.add_paragraph(line)
    document.save(str(output_path))


def _rebuild_xlsx(text: str, output_path: Path) -> None:
    try:
        import openpyxl
    except ImportError as e:
        raise DocumentEngineError("openpyxl não instalado. Rode: pip install openpyxl") from e

    wb = openpyxl.Workbook()
    sheet = wb.active
    sheet.title = "Phoenix"
    for row_idx, line in enumerate(text.split("\n"), 1):
        for col_idx, cell_val in enumerate(line.split(" | "), 1):
            sheet.cell(row=row_idx, column=col_idx, value=cell_val)
    wb.save(str(output_path))


def _rebuild_pptx(text: str, output_path: Path) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError as e:
        raise DocumentEngineError("python-pptx não instalado. Rode: pip install python-pptx") from e

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    slide_chunks = text.split("## Slide")
    if len(slide_chunks) <= 1:
        # LLM não devolveu no formato "## Slide N" esperado - trata o
        # texto inteiro como um único slide em vez de descartar tudo.
        slide_chunks = [text]
    for chunk in slide_chunks:
        chunk = chunk.strip()
        if not chunk:
            continue
        slide = prs.slides.add_slide(blank_layout)
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6))
        textbox.text_frame.text = chunk
    prs.save(str(output_path))
